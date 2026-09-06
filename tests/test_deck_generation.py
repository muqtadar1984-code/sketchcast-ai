"""The slide deck as its OWN generation kind ('deck', 2026-09).

Until now the deck was a by-product of the presentation job, built from the
video script — and the live semantic video prompt dropped the study-note
fields, so decks shipped as heading-plus-glyph. These pin the replacement:

* agent5_slides.deck_notes authors the deck with its own model call and
  normalises the reply into SCRIPT-SHAPED segments, so the slide/deck renderer
  is reused unchanged;
* generate_episode_slides can render into a job's own directory and can skip
  the deck (the presentation's rollout flag);
* worker.process._generate_deck uploads exactly one deck_pptx artifact and
  refuses to finish without a file — here the deck IS the artifact;
* the job rides the fast lane and is a builder for the observer guard.

No model or network call anywhere: the client is a stub, Supabase is the
FakeSB the observer-guard tests use, and the upload is monkeypatched.
"""

from __future__ import annotations

import inspect
import json
import re
import sys
from pathlib import Path

import pytest
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from agent5_slides import deck_notes  # noqa: E402
from agent5_slides.deck_notes import (  # noqa: E402
    DEFAULT_SLIDES, MAX_SLIDES, MIN_SLIDES, author_deck_slides, build_deck_prompt,
    clamp_slides,
)
from agent5_slides.slide_generator import SLIDES_DIR, generate_episode_slides  # noqa: E402
from tests.test_observer_job_guard import FakeSB  # noqa: E402


# ── fixtures ────────────────────────────────────────────────────────────


class _StubClient:
    """Returns a fixed payload (no network) and remembers what it was asked."""

    model = "stub-model"

    def __init__(self, data):
        self._data = data
        self.calls: list[dict] = []
        self.session_usage = {"calls": 1, "input_tokens": 10, "output_tokens": 20, "cost_usd": 0.0}

    def analyze(self, prompt, max_tokens=0, **k):
        self.calls.append({"prompt": prompt, "max_tokens": max_tokens, **k})
        return {"data": self._data, "usage": {}, "truncated": False}


BOOK = {"id": "book-1", "title": "Living Things", "grade": "Year 7", "subject": "Science"}
CHAPTER = {
    "chapter_num": 3,
    "title": "Cells",
    "sections": [{"section_title": "What is a cell?",
                  "content": "All living things are made of cells. The cell membrane controls "
                             "what enters and leaves. The nucleus holds the DNA. Mitochondria "
                             "release energy by respiration.",
                  "subsections": []}],
}
# The Agent-2 shape coverage.measure / chapter_grounding read (concepts and
# episodes are nested one level down, as run_full_analysis().model_dump() emits).
ANALYSIS = {
    "chapter_title": "Cells",
    "concepts": {"concepts": [
        {"concept_id": "c001", "name": "cell membrane"},
        {"concept_id": "c002", "name": "nucleus"},
        {"concept_id": "c003", "name": "mitochondria"},
    ]},
    "episodes": {"episodes": [{"sections_covered": ["What is a cell?"],
                               "key_concepts_introduced": ["c001", "c002", "c003"]}]},
}


def _reply(n_extra_plain: int = 0) -> dict:
    """A realistic model reply: a mix of formats, plus two slides that must be
    dropped (a non-dict and an empty one) and one whose visual cannot render."""
    slides = [
        {"heading": "What are cells made of?", "points": ["The building blocks of life"],
         "visual": None, "notes": "Open with the big question: what is every living thing made of?"},
        {"heading": "Cell membrane", "points": [], "notes": "Introduce the term.",
         "visual": {"kind": "definition", "body": "The thin layer that controls what enters and leaves a cell."}},
        {"heading": "Respiration steps", "points": [], "notes": "Walk the chain left to right.",
         "visual": {"kind": "flow", "nodes": ["Glucose", "Oxygen", "Energy released"], "caption": "In the mitochondria"}},
        {"heading": "Which part holds DNA?", "points": [], "notes": "Quick check.",
         "visual": {"kind": "quiz", "options": ["Membrane", "Nucleus", "Mitochondria"], "answer": 1}},
        # visual kind unknown → dropped by _parse_slide_visual; the points stay
        {"heading": "Parts of a cell", "points": ["Membrane", "Nucleus", "Cytoplasm", "Mitochondria", "Ribosome"],
         "visual": {"kind": "bogus", "nodes": ["a", "b"]}, "notes": "Bullets fallback."},
        # a two-step cycle cannot draw a ring → degrades to a flow
        {"heading": "Energy loop", "points": [], "notes": "Loop.",
         "visual": {"kind": "cycle", "nodes": ["Eat", "Respire"]}},
        "not a slide",
        {"heading": "", "points": [], "visual": None, "notes": "nothing on screen"},
        {"heading": "Takeaways", "points": [], "notes": "Recap the three parts.",
         "visual": {"kind": "takeaways", "nodes": ["Cells are the unit of life", "Membrane controls entry", "Nucleus holds DNA"]}},
    ]
    for i in range(n_extra_plain):
        slides.append({"heading": f"Extra {i}", "points": ["one", "two"], "visual": None, "notes": f"extra {i}"})
    return {"title": "Cells", "slides": slides}


def _segments(n: int = 4) -> list[dict]:
    return [{
        "segment_id": f"d{i:03d}", "type": "explore",
        "text": f"Teacher notes for slide {i}.", "elevenlabs_text": f"Teacher notes for slide {i}.",
        "slide_heading": f"Heading {i}", "slide_points": [f"Point {i}a", f"Point {i}b"],
        "slide_visual": None, "estimated_duration_seconds": 30,
    } for i in range(1, n + 1)]


def _deck_script(segments: list[dict]) -> dict:
    return {"book_id": "book-1", "chapter_num": 3, "episodes": [{
        "book_id": "book-1", "chapter_num": 3, "episode_num": 1,
        "episode_title": "Cells", "segments": segments,
    }]}


# ── (a) the authoring call ──────────────────────────────────────────────


class TestAuthorDeckSlides:
    def test_normalises_the_reply_into_script_shaped_segments(self):
        client = _StubClient(_reply())
        segs = author_deck_slides(BOOK, CHAPTER, ANALYSIS, client, {}, "en")
        # 9 replied: the string and the empty slide are dropped → 7 survive
        assert len(segs) == 7
        assert [s["segment_id"] for s in segs] == [f"d{i:03d}" for i in range(1, 8)]
        for s in segs:
            assert s["type"] == "explore"
            assert s["text"] == s["elevenlabs_text"]
            assert s["estimated_duration_seconds"] == 30
            assert set(s) >= {"segment_id", "type", "text", "elevenlabs_text", "slide_heading",
                              "slide_points", "slide_visual"}
        assert segs[0]["slide_heading"] == "What are cells made of?"
        assert segs[0]["text"].startswith("Open with the big question")
        assert segs[1]["slide_visual"]["kind"] == "definition"
        assert segs[2]["slide_visual"] == {
            "kind": "flow", "nodes": ["Glucose", "Oxygen", "Energy released"], "groups": [],
            "items": [], "caption": "In the mitochondria", "body": "", "options": [], "answer": None,
        }
        assert segs[3]["slide_visual"]["kind"] == "quiz" and segs[3]["slide_visual"]["answer"] == 1

    def test_invalid_visuals_are_dropped_and_points_capped_at_four(self):
        segs = author_deck_slides(BOOK, CHAPTER, ANALYSIS, _StubClient(_reply()), {}, "en")
        bogus = segs[4]
        assert bogus["slide_visual"] is None, "an unknown kind must not reach the renderer"
        assert bogus["slide_points"] == ["Membrane", "Nucleus", "Cytoplasm", "Mitochondria"]
        assert segs[5]["slide_visual"]["kind"] == "flow", "a 2-node cycle degrades to a flow"

    def test_the_call_is_grounded_and_asks_for_the_clamped_count(self):
        client = _StubClient(_reply())
        author_deck_slides(BOOK, CHAPTER, ANALYSIS, client, {"num_slides": 12}, "en")
        call = client.calls[0]
        assert call["max_tokens"] == 8192
        from docgen.docx_builder import chapter_grounding
        assert call["cache_prefix"] == chapter_grounding(BOOK, CHAPTER, ANALYSIS), \
            "byte-identical grounding → the documents' prompt cache is re-read"
        assert "Produce exactly 12 slides" in call["prompt"]
        assert "cell membrane" in call["cache_prefix"]

    @pytest.mark.parametrize("raw, expected", [
        (None, DEFAULT_SLIDES), ("x", DEFAULT_SLIDES), (10, 10), ("7", 7),
        (1, MIN_SLIDES), (0, MIN_SLIDES), (-4, MIN_SLIDES), (16, MAX_SLIDES), (99, MAX_SLIDES),
    ])
    def test_the_slide_count_is_clamped_6_to_16(self, raw, expected):
        assert clamp_slides(raw) == expected
        client = _StubClient(_reply())
        author_deck_slides(BOOK, CHAPTER, ANALYSIS, client, {"num_slides": raw}, "en")
        assert f"Produce exactly {expected} slides" in client.calls[0]["prompt"]

    def test_fewer_than_three_usable_slides_fails_loud(self):
        thin = {"title": "t", "slides": [
            {"heading": "One", "points": ["a"], "notes": "n"},
            {"heading": "Two", "points": ["b"], "notes": "n"},
            "junk", {"heading": "", "points": [], "visual": {"kind": "nope"}},
        ]}
        with pytest.raises(RuntimeError, match="deck authoring returned 2 slides"):
            author_deck_slides(BOOK, CHAPTER, ANALYSIS, _StubClient(thin), {}, "en")

    @pytest.mark.parametrize("data", [{}, None, {"slides": "not a list"}, {"title": "only"}])
    def test_an_empty_or_malformed_reply_fails_loud(self, data):
        with pytest.raises(RuntimeError, match="deck authoring returned 0 slides"):
            author_deck_slides(BOOK, CHAPTER, ANALYSIS, _StubClient(data), {}, "en")

    def test_a_visual_that_names_its_format_under_another_key_is_discarded(self):
        """Why the prompt must spell out the "kind" key (review 2026-09-04):
        the validator is strict on it, so a reply shaped {"type": "flow"} or
        {"flow": {...}} silently ships a heading-only slide — no error,
        MIN_SURVIVING not tripped."""
        reply = {"title": "t", "slides": [
            {"heading": "A", "points": [], "notes": "n", "visual": {"type": "flow", "nodes": ["x", "y"]}},
            {"heading": "B", "points": [], "notes": "n", "visual": {"flow": {"nodes": ["x", "y"]}}},
            {"heading": "C", "points": [], "notes": "n", "visual": {"kind": "flow", "nodes": ["x", "y"]}},
        ]}
        segs = author_deck_slides(BOOK, CHAPTER, ANALYSIS, _StubClient(reply), {}, "en")
        assert [s["slide_visual"] is None for s in segs] == [True, True, False]
        assert segs[0]["slide_points"] == [] and segs[1]["slide_points"] == []


class TestDeckPrompt:
    def test_asks_for_a_deck_with_the_legacy_catalogue(self):
        p = build_deck_prompt(10, "en")
        assert "Produce exactly 10 slides in teaching order" in p
        assert '{"title": "...", "slides": [' in p
        for fmt in ("flow", "cycle", "hierarchy", "compare", "icons", "definition", "quiz", "takeaways"):
            assert f'"{fmt}"' in p
        assert "ANTI-MONOTONY RULE" in p and "caption" in p
        assert "{n}" not in p
        assert "LANGUAGE" not in p, "English carries no directive"

    def test_the_visual_shape_names_the_kind_key_the_validator_requires(self):
        """Review 2026-09-04 (HIGH): the catalogue listed the format names but
        never said which KEY carries them; _parse_slide_visual reads "kind" and
        drops anything else. The prompt must state the key and enumerate every
        valid value in it."""
        from agent3_scripts.script_generator import _VALID_VISUAL_KINDS
        p = build_deck_prompt(10, "en")
        assert '"kind"' in p and '"kind": "flow"' in p
        m = re.search(r'\{"kind": "<([a-z|]+)>"', p)
        assert m, 'the catalogue header must show {"kind": "<flow|cycle|...>"}'
        assert set(m.group(1).split("|")) == _VALID_VISUAL_KINDS
        assert "yours has exactly 10 slides" in p

    def test_the_output_example_is_valid_json_the_validator_accepts(self):
        """The abbreviated example is the shape the model copies: it must parse,
        every visual in it must survive _parse_slide_visual with its kind
        intact, and the four families (diagram / definition / quiz / closer)
        must all be shown, so no format is demonstrated only by prose."""
        from agent3_scripts.script_generator import _parse_slide_visual
        p = deck_notes.PROMPT.replace("{n}", "10")
        example = json.loads(p[p.index('{"title": "...", "slides": ['):].strip())
        assert set(example) == {"title", "slides"} and len(example["slides"]) == 5
        kinds = []
        for s in example["slides"]:
            assert set(s) == {"heading", "points", "visual", "notes"}
            if s["visual"] is None:
                continue
            parsed = _parse_slide_visual(s["visual"])
            assert parsed is not None and parsed.kind == s["visual"]["kind"]
            kinds.append(parsed.kind)
        assert kinds == ["flow", "definition", "quiz", "takeaways"]
        assert example["slides"][0]["visual"] is None, "a plain-points slide shows visual: null"

    def test_the_language_directive_is_appended_like_docgen(self):
        from shared.languages import prompt_directive
        p = build_deck_prompt(8, "ms")
        assert p.endswith(prompt_directive("ms"))
        assert "Bahasa Melayu" in p or "Malay" in p

    def test_jawi_is_two_scripts_on_screen_jawi_notes_rumi(self):
        p = build_deck_prompt(8, "ms-arab")
        assert "JAWI deck" in p and "چ ڠ ڤ ݢ ۏ ڽ" in p
        assert "notes stay Rumi, on-screen stays Jawi" in p
        # NOT the everything-Jawi document directive — that would make the
        # teacher's notes unreadable to a Malay voice or a Rumi-only reader.
        assert "Do NOT output any Rumi/Latin" not in p

    def test_language_reaches_the_authoring_call(self):
        client = _StubClient(_reply())
        author_deck_slides(BOOK, CHAPTER, ANALYSIS, client, {}, "ms-arab")
        assert "JAWI deck" in client.calls[0]["prompt"]


# ── (b)/(c) rendering into the job's own dir ────────────────────────────


class TestGenerateEpisodeSlidesForADeck:
    def test_writes_the_deck_in_out_dir_with_the_notes_per_slide(self, tmp_path):
        segs = _segments(4)
        manifest = generate_episode_slides(script_data=_deck_script(segs), out_dir=tmp_path).model_dump()
        deck = tmp_path / "episode_1_deck.pptx"
        assert manifest["deck_path"] == str(deck) and deck.exists()

        prs = Presentation(str(deck))
        slides = list(prs.slides)
        assert len(slides) == len(segs) + 2, "title + one per segment + closing"
        for seg, slide in zip(segs, slides[1:-1]):
            assert slide.notes_slide.notes_text_frame.text == seg["text"]
            assert any(sh.shape_type is not None and getattr(sh, "image", None) is not None
                       for sh in slide.shapes), "the content slide embeds the rendered PNG"

        # Isolation: every PNG lands under out_dir, not the shared storage dir.
        pngs = sorted(p.name for p in tmp_path.glob("*_slide.png"))
        assert pngs == [f"{s['segment_id']}_slide.png" for s in segs]
        for m in manifest["segments"]:
            assert Path(m["slide_image_path"]).parent == tmp_path
            assert m["slide_path"] == str(deck)
        assert not (SLIDES_DIR / "book-1" / "chapter_3").exists(), \
            "a deck job must never touch storage/slides/{book}/chapter_{n}"

    def test_build_deck_false_renders_pngs_but_no_pptx(self, tmp_path):
        segs = _segments(3)
        manifest = generate_episode_slides(
            script_data=_deck_script(segs), out_dir=tmp_path, build_deck=False,
        ).model_dump()
        assert manifest["deck_path"] is None
        assert list(tmp_path.glob("*.pptx")) == []
        assert len(list(tmp_path.glob("*_slide.png"))) == 3, "the video still gets its slide images"
        assert all(m["slide_path"] is None for m in manifest["segments"])

    def test_the_default_still_builds_the_deck(self, tmp_path):
        sig = inspect.signature(generate_episode_slides)
        assert sig.parameters["build_deck"].default is True
        assert sig.parameters["out_dir"].default is None
        assert sig.parameters["deck_required"].default is False, \
            "the presentation's embedded deck stays a bonus"

    def test_a_deck_build_error_is_swallowed_by_default_but_raised_when_required(self, monkeypatch, tmp_path):
        """Review 2026-09-04 (LOW): the presentation logs a failed deck and
        carries on (bonus semantics, unchanged); a deck JOB gets the ORIGINAL
        error re-raised, so support triage sees the python-pptx/template cause
        rather than the generic 'produced no file'."""
        from agent5_slides import slide_generator as sg

        def _boom(*a, **k):
            raise ValueError("template has no Title Slide layout")

        monkeypatch.setattr(sg, "build_episode_deck", _boom)
        bonus = sg.generate_episode_slides(script_data=_deck_script(_segments(3)),
                                           out_dir=tmp_path / "bonus").model_dump()
        assert bonus["deck_path"] is None and len(list((tmp_path / "bonus").glob("*_slide.png"))) == 3

        with pytest.raises(RuntimeError, match="deck build failed: template has no Title Slide layout") as ei:
            sg.generate_episode_slides(script_data=_deck_script(_segments(3)),
                                       out_dir=tmp_path / "job", deck_required=True)
        assert isinstance(ei.value.__cause__, ValueError)


# ── (d) the worker helper ───────────────────────────────────────────────


def _worker_env(monkeypatch):
    from worker import client as db
    from worker import process

    sb = FakeSB()
    sb.tables["generations"] = [{"id": "gen-1", "status": "processing", "kind": "deck",
                                 "owner_id": "u1", "book_id": "book-1", "params": {}}]
    sb.tables["jobs"] = [{"id": "job-1", "type": "deck", "status": "processing",
                          "generation_id": "gen-1", "progress": 45}]
    uploads: list[tuple[str, str]] = []

    def _upload(_sb, local_path, dest):
        assert Path(local_path).exists() and Path(local_path).stat().st_size > 0
        uploads.append((str(local_path), dest))
        return dest

    monkeypatch.setattr(db, "upload_artifact", _upload)
    return sb, uploads, process


class TestGenerateDeckHelper:
    def test_uploads_exactly_one_deck_pptx_row_at_base_deck_pptx(self, monkeypatch, tmp_path):
        sb, uploads, process = _worker_env(monkeypatch)
        client = _StubClient(_reply())
        title = process._generate_deck(
            sb, "job-1", "gen-1", BOOK, CHAPTER, ANALYSIS, client, {"num_slides": 8},
            {}, "en", "ltr", tmp_path, "u1/gen-1", "Cells",
        )
        rows = sb.tables["artifacts"]
        assert [(r["generation_id"], r["kind"], r["storage_path"]) for r in rows] == \
            [("gen-1", "deck_pptx", "u1/gen-1/deck.pptx")]
        assert [d for _, d in uploads] == ["u1/gen-1/deck.pptx"]
        local, _ = uploads[0]
        assert Path(local).is_relative_to(tmp_path / "deck"), "rendered in the job's own dir"
        assert len(list(Presentation(local).slides)) == 7 + 2
        assert title == "Living Things · Cells · Slide deck"
        assert sb.tables["jobs"][0]["progress"] == 96
        # Coverage is measured and recorded on the generation, never gated.
        cov = (sb.tables["generations"][0].get("params") or {}).get("coverage")
        assert cov and cov[0]["kind"] == "deck" and cov[0]["model"] == "stub-model"

    def test_raises_when_the_builder_returns_no_file(self, monkeypatch, tmp_path):
        sb, uploads, process = _worker_env(monkeypatch)
        from agent5_slides import slide_generator as sg
        from agent5_slides.models import SlideManifest

        def _no_deck(**kw):
            return SlideManifest(manifest_id="m", script_id="s", book_id="book-1", chapter_num=3,
                                 deck_path=None)

        monkeypatch.setattr(sg, "generate_episode_slides", _no_deck)
        with pytest.raises(RuntimeError, match="deck build produced no file"):
            process._generate_deck(
                sb, "job-1", "gen-1", BOOK, CHAPTER, ANALYSIS, _StubClient(_reply()), {},
                {}, "en", "ltr", tmp_path, "u1/gen-1", "Cells",
            )
        assert uploads == [] and sb.tables.get("artifacts", []) == []

    def test_the_builders_own_error_reaches_the_job_error(self, monkeypatch, tmp_path):
        """_generate_deck passes deck_required=True: a build_episode_deck
        exception surfaces with its message (→ jobs.error via finish_job),
        not as the generic 'deck build produced no file'."""
        sb, uploads, process = _worker_env(monkeypatch)
        from agent5_slides import slide_generator as sg

        def _boom(*a, **k):
            raise ValueError("bad branding template: no notes placeholder")

        monkeypatch.setattr(sg, "build_episode_deck", _boom)
        with pytest.raises(RuntimeError, match="deck build failed: bad branding template: no notes placeholder"):
            process._generate_deck(
                sb, "job-1", "gen-1", BOOK, CHAPTER, ANALYSIS, _StubClient(_reply()), {},
                {}, "en", "ltr", tmp_path, "u1/gen-1", "Cells",
            )
        assert uploads == [] and sb.tables.get("artifacts", []) == []
        assert "deck_required=True" in inspect.getsource(process._generate_deck)

    def test_a_thin_authoring_reply_fails_before_anything_is_uploaded(self, monkeypatch, tmp_path):
        sb, uploads, process = _worker_env(monkeypatch)
        with pytest.raises(RuntimeError, match="deck authoring returned 0 slides"):
            process._generate_deck(
                sb, "job-1", "gen-1", BOOK, CHAPTER, ANALYSIS, _StubClient({}), {},
                {}, "en", "ltr", tmp_path, "u1/gen-1", "Cells",
            )
        assert uploads == [] and sb.tables.get("artifacts", []) == []

    def test_the_kind_is_dispatched_after_the_presentation_and_before_the_else(self):
        """The observer/acceptance ordering pin (test_semantic_prompt) needs the
        video artifact recorded before any later branch mentions it; the deck
        branch must therefore sit after the presentation branch."""
        # The kind branches live in _build_from_analysis (the shared build
        # half of process_generation, split out 2026-09-06 for the catalogue).
        from worker.process import _build_from_analysis
        src = inspect.getsource(_build_from_analysis)
        deck = src.index('elif kind == "deck":')
        assert src.index('if kind == "presentation":') < deck
        assert src.index('elif kind == "exam":') < deck
        assert deck < src.index("Unsupported generation kind")
        assert "_generate_deck(" in src[deck:]

    def test_the_presentation_keeps_its_deck_until_the_flag_flips(self, monkeypatch):
        """DECK_IN_PRESENTATION defaults to "1" (rollout step 1: worker first,
        presentations unchanged); "0" passes build_deck=False."""
        from worker.process import _build_from_analysis  # the presentation loop lives here
        src = inspect.getsource(_build_from_analysis)
        assert 'os.getenv("DECK_IN_PRESENTATION", "1").strip() == "1"' in src
        assert "build_deck=_deck_in_presentation" in src


# ── (e) the fast lane ───────────────────────────────────────────────────


def test_the_deck_rides_the_fast_lane():
    from worker import run
    assert "deck" in run.DOC_JOB_TYPES
    assert run.DOC_JOB_TYPES.index("deck") > run.DOC_JOB_TYPES.index("exam")


def test_the_deck_is_a_builder_for_the_observer_guard():
    from worker import client as db
    assert db.generation_to_mirror({"type": "deck", "generation_id": "g"}) == "g"


def test_the_module_reuses_the_video_paths_validation():
    """deck_notes must not grow its own visual validator: the video path's
    _parse_slide_visual is the one contract the renderer was tuned against."""
    from agent3_scripts.script_generator import _parse_slide_visual
    assert deck_notes._parse_slide_visual is _parse_slide_visual
