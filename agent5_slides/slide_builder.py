"""Agent 5: Slide Builder — content slides (PNG for video, PPTX for download).

Each segment becomes a clean 1280x720 lesson slide in the Scholar palette
showing the TEXTBOOK CHAPTER CONTENT (a heading + key bullet points) — not the
narration. The spoken Socratic narration is what plays as audio and what goes
into the PPTX speaker notes.

Rendering uses bundled DejaVu fonts (shipped in ./fonts) so it works on any OS.
"""

from __future__ import annotations

import logging
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

from shared.text_shaping import display_text
from .theme import CANVAS as _BG, GRAPHITE as _MUTED, INK as _INK, LINE as _RULE, TEAL as _GREEN, WHITE as _WHITE

logger = logging.getLogger(__name__)

_WIDTH = 1280
_HEIGHT = 720
# Palette names kept for back-compat: _GREEN is now the Live Ink teal accent and
# _BG the near-white canvas. Values live in theme.py (shared with the deck).

_FONT_DIR = Path(__file__).resolve().parent / "fonts"
_MARGIN_X = 90

# Script-aware font selection. DejaVu covers Latin + Arabic presentation
# forms; Devanagari (Marathi) and Telugu need the bundled Noto families AND
# HarfBuzz shaping (Pillow's RAQM layout engine — enabled at runtime when
# libraqm is present; see nixpacks.toml). Without RAQM those scripts render
# with misplaced vowel signs — we still draw (best-effort) rather than fail.
import re as _re

_DEVA_RE = _re.compile(r"[ऀ-ॿ]")
_TELU_RE = _re.compile(r"[ఀ-౿]")
# Arabic script (Arabic + Persian/Urdu/JAWI in Arabic Extended-A + the
# presentation forms arabic_reshaper emits). Noto Sans Arabic covers the
# Jawi-specific letters (ڤ ڠ ۏ ݢ ڽ چ) that DejaVu lacks.
_ARAB_RE = _re.compile(r"[؀-ۿݐ-ݿࢠ-ࣿﭐ-﷿ﹰ-﻿]")
try:
    from PIL import features as _pil_features

    _RAQM = bool(_pil_features.check("raqm"))
except Exception:  # noqa: BLE001
    _RAQM = False


def _font(bold: bool, size: int, sample: str = "") -> ImageFont.FreeTypeFont:
    if sample and _DEVA_RE.search(sample):
        name = "NotoSansDevanagari-Bold.ttf" if bold else "NotoSansDevanagari-Regular.ttf"
    elif sample and _TELU_RE.search(sample):
        name = "NotoSansTelugu-Bold.ttf" if bold else "NotoSansTelugu-Regular.ttf"
    elif sample and _ARAB_RE.search(sample):
        name = "NotoSansArabic-Bold.ttf" if bold else "NotoSansArabic-Regular.ttf"
    else:
        name = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    # RAQM shapes the Indic scripts (they arrive un-shaped). Arabic/Jawi is
    # PRE-SHAPED by arabic_reshaper into presentation forms before drawing, so
    # it must be drawn WITHOUT RAQM (re-shaping presentation forms breaks them).
    if _RAQM and name.startswith("Noto") and not name.startswith("NotoSansArabic"):
        try:
            return ImageFont.truetype(str(_FONT_DIR / name), size, layout_engine=ImageFont.Layout.RAQM)
        except Exception:  # noqa: BLE001
            pass
    return ImageFont.truetype(str(_FONT_DIR / name), size)


def _wrap(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_w: int) -> list[str]:
    """Greedy word-wrap to a pixel width.

    Wraps on LOGICAL words but measures the SHAPED string — Arabic joins
    letters (narrower than isolated forms), so measuring raw codepoints would
    wrap too early. Returned lines are logical; callers shape at draw time.
    """
    words = text.split()
    lines: list[str] = []
    cur = ""
    for w in words:
        trial = f"{cur} {w}".strip()
        if draw.textlength(display_text(trial), font=font) <= max_w:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


# An element is an ordered list of per-line bounding boxes (x0, y0, x1, y1) to
# reveal left→right; a slide's animation is an ordered list of such elements.
_Box = tuple[int, int, int, int]


def _paste_figure(canvas, draw, region, visual, *, accent, rtl=False) -> list[list[_Box]]:
    """Place a REAL textbook figure (Phase 3) into the content region, framed and
    attributed as a quotation ("From the textbook · p.N") with an optional caption.

    A figure is a raster paste, not vector shapes, so it lives here (compose_slide
    owns the canvas) rather than in diagram_builder. Returns reveal elements, or []
    if the image can't be loaded — the caller then falls back to bullets so a bad
    ``src`` never yields an empty slide.
    """
    src = (visual or {}).get("src")
    if not src or not Path(src).exists():
        return []
    try:
        im = Image.open(src)
        if im.mode in ("RGBA", "LA", "P"):  # flatten transparency onto white
            im = im.convert("RGBA")
            flat = Image.new("RGB", im.size, _WHITE)
            flat.paste(im, mask=im.split()[-1])
            fig = flat
        else:
            fig = im.convert("RGB")
    except Exception:  # noqa: BLE001
        return []

    rx0, ry0, rx1, ry1 = region
    rw, rh = rx1 - rx0, ry1 - ry0
    attribution = str(visual.get("attribution") or "").strip()
    caption = " ".join(str(visual.get("caption") or "").strip().split())

    label_h, pad = 30, 14
    cap_h = 30 if caption else 0
    avail_w = rw - 2 * pad
    avail_h = rh - label_h - cap_h - 2 * pad - 8
    if avail_w < 60 or avail_h < 60:
        return []
    scale = min(avail_w / fig.width, avail_h / fig.height, 1.4)  # cap upscale — don't mush a small figure
    fw, fh = max(1, int(fig.width * scale)), max(1, int(fig.height * scale))
    fig = fig.resize((fw, fh))
    fx = rx0 + (rw - fw) // 2
    fy = ry0 + label_h + pad + max(0, (avail_h - fh) // 2)

    card = [int(fx - pad), int(fy - pad), int(fx + fw + pad), int(fy + fh + pad)]
    draw.rounded_rectangle(card, radius=12, fill=_WHITE, outline=_RULE, width=2)
    canvas.paste(fig, (int(fx), int(fy)))

    # Attribution tab (accent) above the card — the "this is quoted, not ours" signal.
    lab = ("From the textbook" + (f"  ·  {attribution}" if attribution else "")).upper()
    lf = _font(bold=True, size=14)
    lw = int(draw.textlength(lab, font=lf))
    lx, ly = card[0], card[1] - label_h + 2
    draw.rounded_rectangle([lx, ly, lx + lw + 24, ly + 24], radius=6, fill=accent)
    draw.text((lx + 12, ly + 4), lab, fill=_WHITE, font=lf)

    reveal: list[list[_Box]] = [[(lx, ly, lx + lw + 24, ly + 24)], [tuple(card)]]
    if caption:
        cf = _font(bold=False, size=20, sample=caption)
        line = display_text(_wrap(draw, caption, cf, rw)[0], rtl_base=rtl)
        cw = int(draw.textlength(line, font=cf))
        cx = rx0 + (rw - cw) // 2
        cyy = card[3] + 10
        draw.text((cx, cyy), line, fill=_MUTED, font=cf)
        reveal.append([(cx, cyy, cx + cw, cyy + 24)])
    return reveal


def compose_slide(
    heading: str,
    points: list[str],
    footer_text: str = "",
    context_title: str = "",
    fallback_text: str = "",
    accent: tuple[int, int, int] | None = None,
    logo_path: str | None = None,
    visual: dict | None = None,
    number: int | None = None,
    concept: str | None = None,
    direction: str = "ltr",
) -> tuple[Image.Image, list[list[_Box]], list[_Box]]:
    """Render the canonical 1280x720 Live Ink lesson slide + its object layout.

    Returns ``(image, anim_elements, static_boxes)``:

    * ``image`` is the fully-drawn slide (what the Agent 6 native renderer
      animates toward — pixel-identical).
    * ``anim_elements`` is the content in *teaching order* (badge → title →
      bullets / diagram); each element is a list of per-line boxes so the
      renderer can write it on left→right.
    * ``static_boxes`` are regions shown from the first frame — the concept
      illustration (introduced at slide entry, NOT interleaved with the bullet
      reveal), the logo, and any footer.

    ``accent`` overrides the teal accent; ``logo_path`` stamps a logo top-right;
    ``number`` draws the teal number badge (the repeated motif); ``concept`` is
    a keyword/icon name for the right-column illustration.
    """
    acc = accent or _GREEN  # teal
    rtl = direction == "rtl"
    canvas = Image.new("RGB", (_WIDTH, _HEIGHT), _BG)
    draw = ImageDraw.Draw(canvas)
    M = _MARGIN_X
    anim: list[list[_Box]] = []
    static: list[_Box] = []

    # Shape once at draw time (Arabic joins + bidi); no-op for other scripts.
    disp = lambda s: display_text(s, rtl_base=rtl)  # noqa: E731

    # Optional school logo — top-right for LTR, top-LEFT mirrored for RTL.
    if logo_path:
        try:
            logo = Image.open(logo_path).convert("RGBA")
            lh = 52
            lw = max(1, int(logo.width * lh / max(1, logo.height)))
            logo = logo.resize((lw, lh))
            lx, ly = (M, 34) if rtl else (_WIDTH - M - lw, 34)
            canvas.paste(logo, (lx, ly), logo)
            static.append((lx, ly, lx + lw, ly + lh))
        except Exception:  # noqa: BLE001
            pass

    has_visual = bool(visual)
    show_illus = bool(concept) and not has_visual
    # Illustration column: right for LTR, LEFT for RTL (whole layout mirrors).
    content_right = 800 if (show_illus and not rtl) else _WIDTH - M
    content_left = (M + 480) if (show_illus and rtl) else M

    # Number badge (teal circle) — the repeated motif and first reveal element.
    title_x = content_left
    title_right = content_right
    if number is not None:
        br = 27
        bcx, bcy = (_WIDTH - M - br, 92) if rtl else (M + br, 92)
        draw.ellipse([bcx - br, bcy - br, bcx + br, bcy + br], fill=acc)
        nf = _font(bold=True, size=22)
        ns = str(number)
        nw = draw.textlength(ns, font=nf)
        draw.text((bcx - nw / 2, bcy - 15), ns, fill=_WHITE, font=nf)
        anim.append([(bcx - br, bcy - br, bcx + br, bcy + br)])
        if rtl:
            title_right = bcx - br - 22
        else:
            title_x = bcx + br + 22

    # Title (ink), beside the badge — right-anchored for RTL.
    head = (heading or context_title or "SketchCast AI").strip()[:90]
    hf = _font(bold=True, size=38, sample=head)
    ty = 66
    hbox: list[_Box] = []
    for ln in _wrap(draw, head, hf, title_right - title_x)[:2]:
        s = disp(ln)
        lw_px = int(draw.textlength(s, font=hf))
        lx = (title_right - lw_px) if rtl else title_x
        draw.text((lx, ty), s, fill=_INK, font=hf)
        hbox.append((lx, ty, lx + lw_px, ty + 46))
        ty += 50
    if hbox:
        anim.append(hbox)

    # Concept illustration (static — shown at slide entry, opposite the text).
    if show_illus:
        from .theme import draw_concept

        disc = 300
        icx = (M + disc // 2) if rtl else (_WIDTH - M - disc // 2)
        icy = 422
        static.append(draw_concept(canvas, icx, icy, disc, concept, accent=acc))

    y = max(ty + 30, 212)

    # A composable diagram/archetype (or a real textbook figure) claims the content
    # area when present; it falls back to plain bullets if unusable, so slides are
    # never bare.
    diagram_elems: list[list[_Box]] = []
    if has_visual:
        from .diagram_builder import caption_element, render_diagram, set_direction

        set_direction(rtl)  # bidi base for mixed runs inside diagram labels
        region = (content_left, y, content_right, _HEIGHT - 70)
        if isinstance(visual, dict) and visual.get("kind") == "figure":
            # A pasted textbook figure carries its own attribution + caption.
            diagram_elems = _paste_figure(canvas, draw, region, visual, accent=acc, rtl=rtl)
            anim.extend(diagram_elems)
        else:
            diagram_elems = render_diagram(draw, region, visual, accent=acc)
            if diagram_elems:
                anim.extend(diagram_elems)
                cap = caption_element(draw, region, visual.get("caption", "") if isinstance(visual, dict) else "")
                if cap:
                    anim.append(cap)

    if not has_visual or not diagram_elems:
        pts = [p for p in (points or []) if p.strip()]
        text_w = content_right - content_left
        if pts:
            # Bullet points (chapter content) — teal dots + ink text. RTL puts
            # the dot on the RIGHT edge with text right-anchored beside it.
            size = 30
            _pts_sample = " ".join(pts)
            while size >= 20:
                bf = _font(bold=False, size=size, sample=_pts_sample)
                line_h = int(size * 1.5)
                wrapped: list[list[str]] = [_wrap(draw, p, bf, text_w - 40) for p in pts]
                total_h = sum(len(w) * line_h + 16 for w in wrapped)
                if total_h <= (_HEIGHT - y - 60) or size == 20:
                    break
                size -= 2
            dot_r = 7
            for wlines in wrapped:
                ebox: list[_Box] = []
                if rtl:
                    dot_x0 = content_right - 2 * dot_r
                    draw.ellipse([(dot_x0, y + size // 2 - dot_r + 2), (content_right, y + size // 2 + dot_r + 2)], fill=acc)
                else:
                    draw.ellipse([(content_left, y + size // 2 - dot_r + 2), (content_left + 2 * dot_r, y + size // 2 + dot_r + 2)], fill=acc)
                for j, ln in enumerate(wlines):
                    s = disp(ln)
                    lw_px = int(draw.textlength(s, font=bf))
                    if rtl:
                        tx = content_right - 30 - lw_px
                        x1 = content_right if j == 0 else content_right - 30
                        ebox.append((tx, y, x1, y + size))
                    else:
                        tx = content_left + 30
                        x0 = content_left if j == 0 else content_left + 30
                        ebox.append((x0, y, tx + lw_px, y + size))
                    draw.text((tx, y), s, fill=_INK, font=bf)
                    y += line_h
                if ebox:
                    anim.append(ebox)
                y += 16
        # No bullets AND no visual — a hook / question / transition segment. The
        # heading (the hook itself) plus the concept glyph carry the slide; we do
        # NOT print the spoken narration onto it (that read as "the script is on
        # the slide"). fallback_text stays the video/notes voiceover only.

    # Footer is a dev label only — the video composer passes "" in production
    # (gated behind DEBUG_VIDEO), so it never appears in shipped frames.
    if footer_text:
        ff = _font(bold=False, size=16)
        fw = int(draw.textlength(footer_text, font=ff))
        fx = _WIDTH - M - fw
        draw.text((fx, _HEIGHT - 46), footer_text, fill=_MUTED, font=ff)
        static.append((fx, _HEIGHT - 46, _WIDTH - M, _HEIGHT - 46 + 20))

    return canvas, anim, static


def export_slide_png(
    heading: str,
    points: list[str],
    output_png_path: str | Path,
    footer_text: str = "",
    context_title: str = "",
    fallback_text: str = "",
    accent: tuple[int, int, int] | None = None,
    logo_path: str | None = None,
    visual: dict | None = None,
    number: int | None = None,
    concept: str | None = None,
    direction: str = "ltr",
) -> Path:
    """Render a 1280x720 lesson slide PNG: heading + chapter bullets or a diagram.

    Thin wrapper over :func:`compose_slide` (the single source of slide layout,
    shared with the Agent 6 native video renderer) that just saves the image.
    ``number`` (the teal badge) and ``concept`` (the right-column illustration)
    are threaded through so the DECK's slide image matches the animated video
    frame exactly — the same badge and the same space-filling glyph.
    """
    output_png_path = Path(output_png_path)
    output_png_path.parent.mkdir(parents=True, exist_ok=True)

    canvas, _anim, _static = compose_slide(
        heading=heading,
        points=points,
        footer_text=footer_text,
        context_title=context_title,
        fallback_text=fallback_text,
        accent=accent,
        logo_path=logo_path,
        visual=visual,
        number=number,
        concept=concept,
        direction=direction,
    )
    canvas.save(str(output_png_path), "PNG")
    logger.info("Slide PNG: %s (%d points)", output_png_path.name, len([p for p in (points or []) if p.strip()]))
    return output_png_path


# "Live Ink" deck palette — the same single source as the video (theme.py), so
# deck and video can't drift: one dominant ink, white content, one teal accent.
from .theme import (  # noqa: E402
    FAINT as _LI_DIM,
    INK as _LI_INK,
    TEAL as _LI_TEAL,
    WHITE as _LI_WHITE,
)


def build_episode_deck(
    slides: list[dict],
    output_pptx_path: str | Path,
    episode_title: str = "",
    template: str | None = None,
    accent: tuple[int, int, int] | None = None,
    direction: str = "ltr",
) -> Path:
    """Build one editable PPTX deck: a designed lesson deck, narration in notes.

    ``slides`` is a list of {"heading", "points", "narration"} dicts. Without a
    ``template`` it builds the designed "Live Ink" deck (dark title/closing,
    light content, icon-in-circle motif, alternating layouts). When a school
    ``template`` (.pptx) is given the deck inherits that theme instead; ``accent``
    overrides the brand colour in both.
    """
    output_pptx_path = Path(output_pptx_path)
    output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    if template:
        path = _build_branded_deck(slides, output_pptx_path, episode_title, template, accent)
    else:
        path = _build_designed_deck(slides, output_pptx_path, episode_title, accent)
    # RTL (Arabic) lessons: PowerPoint does the SHAPING itself — the deck only
    # needs paragraph-level right alignment + the rtl flag. Applied as a post-
    # pass over the saved file so both deck builders stay untouched.
    if direction == "rtl":
        try:
            _mirror_deck_rtl(path)
        except Exception as exc:  # noqa: BLE001
            logger.warning("RTL deck pass failed for %s: %s", path, exc)
    return path


def _mirror_deck_rtl(pptx_path: Path) -> None:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for p in shape.text_frame.paragraphs:
                # Centred paragraphs (number badges, title slides) stay centred.
                if p.alignment != PP_ALIGN.CENTER:
                    p.alignment = PP_ALIGN.RIGHT
                # python-pptx has no first-class rtl API — set the OOXML attr.
                p._p.get_or_add_pPr().set("rtl", "1")  # noqa: SLF001
    prs.save(str(pptx_path))


def _build_branded_deck(slides, output_pptx_path, episode_title, template, accent):
    """School-template path — inherit the uploaded theme (unchanged behaviour)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    prs = Presentation(template)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    head_rgb = RGBColor(*(accent or _GREEN))
    for s in slides:
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Emu(700000), Emu(900000), Emu(10800000), Emu(5000000))
        tf = box.text_frame
        tf.word_wrap = True
        head = tf.paragraphs[0]
        head.text = (s.get("heading") or episode_title or "SketchCast AI")[:90]
        head.font.size = Pt(30)
        head.font.bold = True
        head.font.color.rgb = head_rgb
        for pt in s.get("points") or []:
            para = tf.add_paragraph()
            para.text = f"•  {pt}"
            para.font.size = Pt(20)
        slide.notes_slide.notes_text_frame.text = s.get("narration") or ""
    prs.save(str(output_pptx_path))
    logger.info("Episode deck (branded) saved: %s (%d content slides)", output_pptx_path.name, len(slides))
    return output_pptx_path


def _build_designed_deck(slides, output_pptx_path, episode_title, accent):
    """Default path — the content slides ARE the rendered lesson slides.

    Each content slide embeds the exact 1280x720 image the video animates toward
    (``compose_slide`` output — the designed layout, diagrams and all),
    full-bleed, with the Socratic narration in the speaker notes. So the
    downloadable deck finally matches what the lesson shows, instead of a heading
    plus flattened bullet labels. A dark title and closing slide bookend it.

    Print sharpness note: the embedded image is 1280x720 (~96 DPI across a 13.3in
    slide) — crisp projected and on-screen, softer in print. A true higher-res
    render means scaling the shared ``compose_slide``/``diagram_builder`` geometry
    (which the video path also uses), so it's deliberately deferred, not bundled
    into this parity pass.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Pt

    IN = 914400
    INK, WHITE = RGBColor(*_LI_INK), RGBColor(*_LI_WHITE)
    ACC, DIM = RGBColor(*(accent or _LI_TEAL)), RGBColor(*_LI_DIM)
    BODY = "Calibri"
    SW, SH = Emu(12192000), Emu(6858000)  # 16:9 canvas — same aspect as the 1280x720 image

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    def new_slide(bg):
        s = prs.slides.add_slide(blank)
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = bg
        return s

    def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        tf = s.shapes.add_textbox(int(l), int(t), int(w), int(h)).text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        return tf

    def para(tf, text, size, color, bold=False, first=False, space=6):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = BODY
        r.font.color.rgb = color
        return p

    def circle(s, cx, cy, d, fill, label="", label_size=15):
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - d / 2), int(cy - d / 2), int(d), int(d))
        o.fill.solid()
        o.fill.fore_color.rgb = fill
        o.line.fill.background()
        o.shadow.inherit = False
        if label:
            tf = o.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = label
            r.font.size = Pt(label_size)
            r.font.bold = True
            r.font.name = BODY
            r.font.color.rgb = WHITE
        return o

    # Title (dark)
    s = new_slide(INK)
    circle(s, 1.0 * IN, 0.95 * IN, 0.62 * IN, ACC, "S", label_size=20)
    tf = textbox(s, 1.0 * IN, 2.45 * IN, 11.3 * IN, 2.3 * IN)
    para(tf, episode_title or "SketchCast Lesson", 40, WHITE, bold=True, first=True, space=8)
    para(tf, "A narrated lesson", 18, ACC, space=0)
    para(textbox(s, 1.0 * IN, 6.7 * IN, 11.3 * IN, 0.5 * IN), "SketchCast AI", 12, DIM, first=True, space=0)

    # Content — one slide per segment, the rendered lesson image full-bleed.
    for i, sd in enumerate(slides, 1):
        s = new_slide(WHITE)
        img = sd.get("image")
        if img and Path(img).exists():
            s.shapes.add_picture(str(img), 0, 0, width=SW, height=SH)
        else:
            # A render slipped — never ship a blank slide: fall back to the
            # heading + bullets so the content still lands (rare).
            heading = (sd.get("heading") or episode_title or "Lesson")[:90]
            circle(s, 0.95 * IN, 1.02 * IN, 0.52 * IN, ACC, str(i))
            para(textbox(s, 1.45 * IN, 0.7 * IN, 11.1 * IN, 0.95 * IN, anchor=MSO_ANCHOR.MIDDLE),
                 heading, 27, INK, bold=True, first=True, space=0)
            y = 2.4 * IN
            for pt in [str(p).strip() for p in (sd.get("points") or []) if str(p).strip()][:6]:
                circle(s, 1.04 * IN, y + 0.13 * IN, 0.17 * IN, ACC)
                para(textbox(s, 1.31 * IN, y - 0.03 * IN, 10.7 * IN, 0.7 * IN), pt, 16, INK, first=True, space=0)
                y += 0.7 * IN
        s.notes_slide.notes_text_frame.text = sd.get("narration") or ""

    # Closing (dark)
    s = new_slide(INK)
    circle(s, 1.0 * IN, 0.95 * IN, 0.62 * IN, ACC)
    tf = textbox(s, 1.0 * IN, 2.5 * IN, 11.3 * IN, 2.4 * IN)
    para(tf, "Ready to teach.", 34, WHITE, bold=True, first=True, space=8)
    para(tf, episode_title or "", 18, ACC, space=12)
    para(tf, "Every lesson slide's speaker notes carry the narration.", 14, DIM, space=0)
    para(textbox(s, 1.0 * IN, 6.7 * IN, 11.3 * IN, 0.5 * IN), "SketchCast AI", 12, DIM, first=True, space=0)

    prs.save(str(output_pptx_path))
    logger.info("Episode deck (designed, image slides) saved: %s (%d content slides)", output_pptx_path.name, len(slides))
    return output_pptx_path
