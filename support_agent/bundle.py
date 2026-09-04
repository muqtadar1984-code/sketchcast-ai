"""Diagnostic-bundle assembly — tenant-scoped by construction.

Every read in this module is keyed to ids taken from the issue row itself
(generation_id / book_id) after ``assert_scope`` proves the reporter may see
that content. Nothing here accepts a free query; there is no code path that
touches another tenant's rows.
"""

from __future__ import annotations

import logging
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

_SOURCE_SNIPPET_CHARS = 4000
_ARTIFACT_SNIPPET_CHARS = 4000


class ScopeViolation(Exception):
    """The reporter has no right to the content named on the issue."""


def assert_scope(sb, reporter_id: str, gen: dict | None, book: dict | None) -> None:
    """EVERY referenced object must be in the reporter's scope INDEPENDENTLY —
    the generation AND the book. Never OR-collapse: a generation the reporter
    owns can carry a foreign book_id (the FK only requires existence), and
    collapsing to "owns the generation" would hand the service-role worker a
    cross-tenant read/write primitive. Raises ScopeViolation otherwise."""
    if not reporter_id or (gen is None and book is None):
        raise ScopeViolation("issue references no owned content")

    rep = (
        sb.table("profiles").select("role, school_id").eq("id", reporter_id).maybe_single().execute()
    )
    rep_d = getattr(rep, "data", None) or {}
    if rep_d.get("role") == "student":
        raise ScopeViolation("student accounts cannot invoke the support agent")

    def _owner_ok(owner_id: str | None, what: str) -> None:
        if not owner_id:
            raise ScopeViolation(f"{what} has no owner")
        if owner_id == reporter_id:
            return
        own = sb.table("profiles").select("school_id").eq("id", owner_id).maybe_single().execute()
        own_d = getattr(own, "data", None) or {}
        if (
            rep_d.get("role") == "school_admin"
            and rep_d.get("school_id")
            and rep_d.get("school_id") == own_d.get("school_id")
        ):
            return
        raise ScopeViolation(f"reporter is outside the {what}'s tenant")

    if gen is not None:
        _owner_ok(gen.get("owner_id"), "generation")
    if book is not None:
        # Book scope: reporter owns it, or it belongs to the reporter's school
        # (the school library case), or the admin rule above via its owner.
        b_owner = book.get("owner_id")
        b_school = book.get("school_id")
        if b_owner == reporter_id:
            pass
        elif b_school and rep_d.get("school_id") and b_school == rep_d.get("school_id"):
            pass
        else:
            _owner_ok(b_owner, "book")
    # Cross-consistency: a generation referencing a book outside its own
    # owner's reach is malformed regardless of the reporter.
    if gen is not None and book is not None:
        if book.get("owner_id") != gen.get("owner_id") and (
            not book.get("school_id") or book.get("school_id") != _school_of(sb, gen.get("owner_id"))
        ):
            raise ScopeViolation("generation references a book outside its owner's tenant")


def _school_of(sb, user_id: str | None) -> str | None:
    if not user_id:
        return None
    r = sb.table("profiles").select("school_id").eq("id", user_id).maybe_single().execute()
    return (getattr(r, "data", None) or {}).get("school_id")


def _docx_text(path: Path, limit: int) -> str:
    try:
        import docx  # python-docx (already a docgen dependency)

        d = docx.Document(str(path))
        return " ".join(p.text for p in d.paragraphs if p.text)[:limit]
    except Exception as exc:  # noqa: BLE001
        logger.warning("docx text extraction failed: %s", exc)
        return ""


def _pptx_text(path: Path, limit: int) -> str:
    try:
        from pptx import Presentation

        out: list[str] = []
        for slide in Presentation(str(path)).slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    out.append(shape.text_frame.text)
        return " ".join(out)[:limit]
    except Exception as exc:  # noqa: BLE001
        logger.warning("pptx text extraction failed: %s", exc)
        return ""


def _chapter_source_text(sb, book: dict, chapter_ref: str | None, tmp: Path) -> tuple[str, dict]:
    """The requested chapter's source text (first N chars) + its stored def.
    Heuristics only (client=None) — this must stay cheap; the model judges."""
    from worker import client as db
    from agent1_ingestion.extractor import extract_pdf
    from agent1_ingestion.structurer import structure_book

    chapters = book.get("chapters") or []
    want = None
    try:
        want = int(str(chapter_ref).strip()) if chapter_ref is not None else None
    except ValueError:
        pass
    cdef = next((c for c in chapters if int(c.get("num", -1)) == want), None)

    pdf_path = db.download_book(sb, book["storage_path"], tmp / "book.pdf")
    extraction = extract_pdf(str(pdf_path))
    structured = structure_book(
        book_id=book["id"], title=book.get("title") or "?", author="?", isbn=None,
        extraction=extraction, images=[], pdf_path=str(pdf_path), client=None,
        known_chapters=chapters or None,
    ).model_dump()
    pick = next(
        (c for c in structured.get("chapters", []) if int(c.get("chapter_num", -1)) == (want or 0)),
        (structured.get("chapters") or [{}])[0],
    )
    text = " ".join(
        (s.get("content") or "") + " "
        + " ".join(ss.get("content") or "" for ss in (s.get("subsections") or []))
        for s in (pick.get("sections") or [])
    )
    meta = {
        "requested_chapter_ref": chapter_ref,
        "stored_chapter": {k: cdef.get(k) for k in ("num", "title", "start_page", "end_page")} if cdef else None,
        "detected_chapter_count": len(chapters),
        "total_pages": extraction.total_pages,
        "has_text_layer": sum(len(i.text or "") for i in extraction.items) >= 200,
    }
    return " ".join(text.split())[:_SOURCE_SNIPPET_CHARS], meta


def assemble_bundle(sb, issue: dict) -> dict:
    """Everything the model needs, all keyed to the issue's own content."""
    reporter_id = issue.get("reporter_id")

    gen = None
    if issue.get("generation_id"):
        r = sb.table("generations").select("*").eq("id", issue["generation_id"]).maybe_single().execute()
        gen = getattr(r, "data", None)
    book = None
    book_id = (gen or {}).get("book_id") or issue.get("book_id")
    if book_id:
        r = sb.table("books").select("*").eq("id", book_id).maybe_single().execute()
        book = getattr(r, "data", None)

    assert_scope(sb, reporter_id, gen, book)

    bundle: dict = {
        "issue": {k: issue.get(k) for k in ("id", "category", "title", "description", "trigger_source")},
        "generation": {k: (gen or {}).get(k) for k in ("id", "kind", "chapter_ref", "status", "title")} if gen else None,
        "book": {k: (book or {}).get(k) for k in ("id", "title", "status")} if book else None,
        "chapters": [
            {k: c.get(k) for k in ("num", "title", "start_page", "end_page")}
            for c in ((book or {}).get("chapters") or [])
        ][:40],
    }

    # Last job for this generation: error text + spend.
    if gen:
        r = (
            sb.table("jobs")
            .select("id, status, error, usage, created_at")
            .eq("generation_id", gen["id"])
            .order("created_at", desc=True)
            .limit(3)
            .execute()
        )
        jobs = getattr(r, "data", None) or []
        bundle["recent_jobs"] = [
            {"status": j.get("status"), "error": (j.get("error") or "")[:300], "usage": j.get("usage")}
            for j in jobs
        ]
        # Is the artifact currently in front of students? (drives the guard)
        r = sb.table("generation_shares").select("id").eq("generation_id", gen["id"]).limit(1).execute()
        bundle["assigned_to_students"] = bool(getattr(r, "data", None))

    # Source + artifact text (wrong-content diagnostics). Best-effort.
    if book and book.get("storage_path"):
        try:
            with tempfile.TemporaryDirectory() as tmp:
                text, meta = _chapter_source_text(sb, book, (gen or {}).get("chapter_ref"), Path(tmp))
                bundle["source_text"] = text
                bundle["source_meta"] = meta
        except Exception as exc:  # noqa: BLE001
            logger.warning("source text unavailable: %s", exc)
            bundle["source_text"] = ""
            bundle["source_meta"] = {"error": str(exc)[:200]}

    if gen:
        try:
            r = sb.table("artifacts").select("kind, storage_path").eq("generation_id", gen["id"]).execute()
            arts = getattr(r, "data", None) or []
            art = next((a for a in arts if a["kind"] == "docx"), None) or next(
                (a for a in arts if a["kind"] == "deck_pptx"), None
            )
            if art is None and gen.get("kind") == "presentation":
                # Since the deck became a generation of its own (DECK_IN_PRESENTATION=0),
                # a lesson carries no deck_pptx row; its text lives on the sibling
                # 'deck' generation of the same unit.
                art = _sibling_deck_artifact(sb, gen)
            if art:
                from worker import client as db

                with tempfile.TemporaryDirectory() as tmp:
                    local = db.download_artifact(sb, art["storage_path"], Path(tmp) / Path(art["storage_path"]).name) \
                        if hasattr(db, "download_artifact") else None
                    if local is None:
                        # storage download via the same bucket helper the worker uses
                        data = sb.storage.from_("artifacts").download(art["storage_path"])
                        local = Path(tmp) / Path(art["storage_path"]).name
                        local.write_bytes(data)
                    bundle["artifact_kind"] = art["kind"]
                    bundle["artifact_text"] = (
                        _docx_text(local, _ARTIFACT_SNIPPET_CHARS)
                        if art["kind"] == "docx"
                        else _pptx_text(local, _ARTIFACT_SNIPPET_CHARS)
                    )
        except Exception as exc:  # noqa: BLE001
            logger.warning("artifact text unavailable: %s", exc)
            bundle["artifact_text"] = ""

    return bundle


def _unit_of(gen: dict) -> tuple:
    """(book, chapter, part) — the unit a kit's generations share."""
    params = gen.get("params") if isinstance(gen.get("params"), dict) else {}
    return (gen.get("book_id"), gen.get("chapter_ref"), str(params.get("part") or ""))


def _sibling_deck_artifact(sb, gen: dict) -> dict | None:
    """The deck_pptx row of the finished 'deck' generation for the same unit as
    a presentation, newest first. Best-effort: None when there is none or the
    lookup fails (the bundle then has no artifact text, as before)."""
    try:
        r = (sb.table("generations").select("id, book_id, chapter_ref, params, created_at")
             .eq("owner_id", gen.get("owner_id")).eq("kind", "deck").eq("status", "done").execute())
        rows = [row for row in (getattr(r, "data", None) or []) if _unit_of(row) == _unit_of(gen)]
        rows.sort(key=lambda row: str(row.get("created_at") or ""), reverse=True)
        for row in rows:
            a = sb.table("artifacts").select("kind, storage_path").eq("generation_id", row["id"]).execute()
            art = next((x for x in (getattr(a, "data", None) or []) if x.get("kind") == "deck_pptx"), None)
            if art:
                return art
    except Exception as exc:  # noqa: BLE001
        logger.warning("sibling deck lookup failed: %s", exc)
    return None
